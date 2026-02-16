import fitz  # PyMuPDF
import pandas as pd
from pptx import Presentation
from docx import Document
import os
from backend.config import settings
from backend.logger import setup_logger

logger = setup_logger(__name__)

def load_document(file_path: str) -> str:
    logger.info(f"Loading document: {file_path}")
    ext = os.path.splitext(file_path)[1].lower()
    text = ""
    
    try:
        if ext == ".pdf":
            doc = fitz.open(file_path)
            for page in doc:
                text += page.get_text()
        elif ext in [".xlsx", ".xls"]:
            df = pd.read_excel(file_path)
            text = df.to_string()
        elif ext == ".csv":
            df = pd.read_csv(file_path)
            text = df.to_string()
        elif ext in [".docx", ".doc"]:
            doc = Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
        elif ext in [".pptx", ".ppt"]:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif ext == ".txt":
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()
        else:
            raise ValueError(f"Unsupported file format: {ext}")
        
        logger.info(f"Loaded {len(text)} characters from {file_path}")
        return text
    except Exception as e:
        logger.error(f"Error loading {file_path}: {e}")
        return ""

def split_text(text: str) -> list[str]:
    chunk_size = settings.CHUNK_SIZE
    overlap = settings.CHUNK_OVERLAP
    logger.info(f"Splitting text with chunk_size={chunk_size}, overlap={overlap}")
    
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start += chunk_size - overlap
    
    logger.info(f"Created {len(chunks)} chunks")
    return chunks
